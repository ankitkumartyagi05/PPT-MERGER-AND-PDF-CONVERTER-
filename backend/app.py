"""
PPTX Merger & Converter — Backend
FastAPI application for uploading, reordering, merging,
and converting PowerPoint presentations to PDF.
"""

import os
import sys
import copy
import uuid
import shutil
import logging
import subprocess
from io import BytesIO
from pathlib import Path
from typing import List, Optional
from datetime import datetime, timedelta

from fastapi import (
    FastAPI, UploadFile, File, HTTPException,
    Form, Query, Request
)
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ──────────────────────────────────────────────────────────────────
# Configuration
# ──────────────────────────────────────────────────────────────────

BASE_DIR = Path(__file__).parent.resolve()
UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "outputs"
SESSIONS_DIR = BASE_DIR / "sessions"

MAX_FILE_SIZE = 100 * 1024 * 1024      # 100 MB per file
MAX_TOTAL_SIZE = 500 * 1024 * 1024     # 500 MB total per session
MAX_FILES = 50                          # Max files per session
ALLOWED_EXTENSIONS = {".pptx"}
SESSION_TTL_HOURS = 2

# ──────────────────────────────────────────────────────────────────
# Logging
# ──────────────────────────────────────────────────────────────────

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)],
)
logger = logging.getLogger(__name__)

# ──────────────────────────────────────────────────────────────────
# FastAPI App
# ──────────────────────────────────────────────────────────────────

app = FastAPI(
    title="PPTX Merger & Converter",
    version="1.0.0",
    description="Upload, reorder, merge, and convert PowerPoint presentations.",
)

frontend_origin = os.getenv("FRONTEND_ORIGIN", "*").strip() or "*"
allow_origins = [frontend_origin] if frontend_origin != "*" else ["*"]

app.add_middleware(
    CORSMiddleware,
    allow_origins=allow_origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Ensure working directories exist
for _dir in [UPLOAD_DIR, OUTPUT_DIR, SESSIONS_DIR]:
    _dir.mkdir(parents=True, exist_ok=True)

# ──────────────────────────────────────────────────────────────────
# Utility helpers
# ──────────────────────────────────────────────────────────────────

def is_valid_pptx(filename: str, content: bytes) -> tuple[bool, str]:
    """Validate file extension and ZIP magic bytes."""
    ext = Path(filename).suffix.lower()
    if ext != ".pptx":
        return False, f"Invalid extension '{ext}'. Only .pptx is accepted."
    if len(content) < 4:
        return False, "File is too small to be a valid PPTX."
    if content[:2] != b"PK":
        return False, "File is not a valid PPTX (not a ZIP archive)."
    return True, ""


def generate_session_id() -> str:
    return uuid.uuid4().hex[:16]


def get_session_path(sid: str) -> Path:
    return SESSIONS_DIR / sid


def format_size(n: int) -> str:
    size = float(n)
    for unit in ("B", "KB", "MB", "GB"):
        if size < 1024:
            return f"{size:.1f} {unit}"
        size /= 1024
    return f"{size:.1f} TB"


def count_slides(path: Path) -> Optional[int]:
    try:
        return len(Presentation(str(path)).slides)
    except Exception:
        return None


def cleanup_stale():
    """Remove sessions and outputs older than SESSION_TTL_HOURS."""
    cutoff = datetime.now() - timedelta(hours=SESSION_TTL_HOURS)
    for d in SESSIONS_DIR.iterdir():
        if d.is_dir() and datetime.fromtimestamp(d.stat().st_mtime) < cutoff:
            shutil.rmtree(d, ignore_errors=True)
            logger.info("Cleaned session %s", d.name)
    for f in OUTPUT_DIR.iterdir():
        if f.is_file() and datetime.fromtimestamp(f.stat().st_mtime) < cutoff:
            f.unlink(missing_ok=True)


# ──────────────────────────────────────────────────────────────────
# Slide merging engine
# ──────────────────────────────────────────────────────────────────

def _copy_slide_shapes(src_slide, dst_slide, src_pres, dst_pres):
    """
    Copy every shape from *src_slide* to *dst_slide*.

    Pictures are re-inserted via ``add_picture`` so that the image
    binary and relationship are correctly stored in the destination
    package.  Everything else (text boxes, arrows, tables, etc.) is
    deep-copied at the XML level.
    """
    # Strip any placeholder shapes the blank layout may have added
    for shape in list(dst_slide.shapes):
        shape._element.getparent().remove(shape._element)

    for shape in src_slide.shapes:
        try:
            # ── Picture: re-insert with proper relationship ──
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                dst_slide.shapes.add_picture(
                    BytesIO(img.blob),
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height,
                )
                continue

            # ── Everything else: XML deep-copy ──
            new_el = copy.deepcopy(shape._element)
            dst_slide.shapes._spTree.append(new_el)

        except Exception as exc:
            logger.warning("Shape %s copy failed: %s", shape.shape_id, exc)


def merge_pptx_files(paths: List[Path], output: Path) -> dict:
    """Merge many PPTX files into one, returning stats + warnings."""
    if not paths:
        raise ValueError("No files to merge.")

    warnings: list[str] = []
    total_slides = 0

    # Use the first file as the base presentation
    try:
        merged = Presentation(str(paths[0]))
        ref_size = (merged.slide_width, merged.slide_height)
        total_slides = len(merged.slides)
    except Exception as exc:
        raise RuntimeError(f"Cannot open '{paths[0].name}': {exc}") from exc

    for fpath in paths[1:]:
        try:
            src = Presentation(str(fpath))
            if (src.slide_width, src.slide_height) != ref_size:
                warnings.append(
                    f"'{fpath.name}' has different slide dimensions — "
                    "content may not align perfectly."
                )
            added = 0
            for slide in src.slides:
                try:
                    blank = (
                        merged.slide_layouts[6]
                        if len(merged.slide_layouts) > 6
                        else merged.slide_layouts[0]
                    )
                    ns = merged.slides.add_slide(blank)
                    _copy_slide_shapes(slide, ns, src, merged)
                    added += 1
                except Exception as exc:
                    warnings.append(
                        f"In '{fpath.name}': slide skipped ({str(exc)[:100]})"
                    )
            total_slides += added
        except Exception as exc:
            warnings.append(f"'{fpath.name}' skipped: {str(exc)[:100]}")
            logger.error("Skipping %s: %s", fpath.name, exc)

    if total_slides == 0:
        raise RuntimeError("No slides could be merged from any file.")

    merged.save(str(output))
    logger.info("Merged → %s  (%d slides)", output.name, total_slides)
    return {"total_slides": total_slides, "warnings": warnings[:50]}


# ──────────────────────────────────────────────────────────────────
# PDF conversion engine
# ──────────────────────────────────────────────────────────────────

def _convert_via_libreoffice(inp: Path, out: Path) -> bool:
    """Attempt conversion with LibreOffice (cross-platform)."""
    candidates = ["libreoffice", "soffice"]
    if sys.platform == "win32":
        for p in (
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ):
            if Path(p).exists():
                candidates.insert(0, p)
    elif sys.platform == "darwin":
        mac = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if Path(mac).exists():
            candidates.insert(0, mac)

    for cmd_base in candidates:
        try:
            logger.info("Trying %s …", cmd_base)
            proc = subprocess.run(
                [cmd_base, "--headless", "--convert-to", "pdf",
                 "--outdir", str(out.parent), str(inp)],
                capture_output=True, timeout=180, cwd=str(out.parent),
            )
            if proc.returncode == 0:
                generated = out.parent / (inp.stem + ".pdf")
                if generated.exists():
                    if generated.resolve() != out.resolve():
                        shutil.move(str(generated), str(out))
                    return True
            logger.warning(
                "%s exited %d: %s",
                cmd_base, proc.returncode,
                proc.stderr.decode(errors="replace")[:300],
            )
        except FileNotFoundError:
            continue
        except subprocess.TimeoutExpired:
            logger.warning("%s timed out", cmd_base)
        except Exception as exc:
            logger.warning("%s error: %s", cmd_base, exc)
    return False


def _convert_via_com(inp: Path, out: Path) -> bool:
    """Attempt conversion via PowerPoint COM (Windows only)."""
    if sys.platform != "win32":
        return False
    try:
        import win32com.client  # type: ignore
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        ppt.Visible = False
        ppt.DisplayAlerts = False
        pres = ppt.Presentations.Open(str(inp.resolve()), WithWindow=False)
        pres.SaveAs(str(out.resolve()), 32)  # ppSaveAsPDF
        pres.Close()
        ppt.Quit()
        return out.exists()
    except ImportError:
        return False
    except Exception as exc:
        logger.error("COM automation failed: %s", exc)
        return False


def convert_to_pdf(inp: Path, out: Path) -> dict:
    """Try LibreOffice, then COM. Return dict with success/method/error."""
    if _convert_via_libreoffice(inp, out):
        return {"success": True, "method": "LibreOffice"}
    if _convert_via_com(inp, out):
        return {"success": True, "method": "PowerPoint COM"}
    return {
        "success": False,
        "method": None,
        "error": (
            "No PDF converter found. Install LibreOffice "
            "(libreoffice.org) or Microsoft PowerPoint."
        ),
    }


# ──────────────────────────────────────────────────────────────────
# API Endpoints
# ──────────────────────────────────────────────────────────────────

@app.get("/", response_class=HTMLResponse)
async def serve_frontend():
    """Serve the single-page frontend."""
    html = BASE_DIR / "index.html"
    if not html.exists():
        return HTMLResponse("<h1>index.html not found</h1>", status_code=500)
    return HTMLResponse(html.read_text(encoding="utf-8"))


@app.post("/upload")
async def upload_files(
    files: List[UploadFile] = File(...),
    session_id: Optional[str] = Form(None),
):
    """Upload .pptx files (new session or append to existing)."""
    if not files:
        raise HTTPException(400, "No files provided.")
    if len(files) > MAX_FILES:
        raise HTTPException(400, f"Maximum {MAX_FILES} files per session.")

    # Resolve session
    if session_id:
        sp = get_session_path(session_id)
        if not sp.exists():
            raise HTTPException(404, "Session not found or expired.")
    else:
        session_id = generate_session_id()
        sp = get_session_path(session_id)
        sp.mkdir(parents=True, exist_ok=True)

    existing = sorted(sp.glob("*"))
    start_idx = len(existing)
    uploaded = []
    total_size = sum(f.stat().st_size for f in existing)

    try:
        for i, uf in enumerate(files):
            fname = uf.filename or f"file_{i}.pptx"
            raw = await uf.read()
            fsize = len(raw)

            ok, msg = is_valid_pptx(fname, raw)
            if not ok:
                raise HTTPException(400, f"{fname}: {msg}")
            if fsize > MAX_FILE_SIZE:
                raise HTTPException(
                    400,
                    f"{fname} exceeds {format_size(MAX_FILE_SIZE)} limit.",
                )
            total_size += fsize
            if total_size > MAX_TOTAL_SIZE:
                raise HTTPException(
                    400,
                    f"Total upload exceeds {format_size(MAX_TOTAL_SIZE)}.",
                )

            safe = f"{start_idx + i:04d}_{Path(fname).name}"
            dest = sp / safe
            dest.write_bytes(raw)

            slides = count_slides(dest)
            uploaded.append(
                {
                    "id": safe,
                    "original_name": fname,
                    "size": fsize,
                    "size_display": format_size(fsize),
                    "slide_count": slides,
                }
            )
            await uf.seek(0)
    except HTTPException:
        if start_idx == 0:
            shutil.rmtree(sp, ignore_errors=True)
        raise
    except Exception as exc:
        if start_idx == 0:
            shutil.rmtree(sp, ignore_errors=True)
        raise HTTPException(500, f"Upload error: {exc}") from exc

    return {"session_id": session_id, "files": uploaded}


@app.post("/remove-file")
async def remove_file(
    session_id: str = Form(...),
    file_id: str = Form(...),
):
    """Delete a single uploaded file from a session."""
    sp = get_session_path(session_id)
    if not sp.exists():
        raise HTTPException(404, "Session not found.")
    target = sp / file_id
    if not target.exists():
        raise HTTPException(404, "File not found.")
    target.unlink()
    return {"status": "ok"}


@app.post("/process")
async def process_merge(
    session_id: str = Form(...),
    file_order: str = Form(...),
):
    """Merge uploaded files in the given order."""
    sp = get_session_path(session_id)
    if not sp.exists():
        raise HTTPException(404, "Session not found or expired.")

    order = [f.strip() for f in file_order.split(",") if f.strip()]
    if not order:
        raise HTTPException(400, "No files specified for merging.")

    paths = []
    for fid in order:
        p = sp / fid
        if not p.exists():
            raise HTTPException(400, f"File not found: {fid}")
        paths.append(p)

    out_name = f"merged_{session_id}.pptx"
    out_path = OUTPUT_DIR / out_name

    try:
        result = merge_pptx_files(paths, out_path)
    except Exception as exc:
        raise HTTPException(500, str(exc)) from exc

    return {
        "status": "success",
        "output_file": out_name,
        "total_slides": result["total_slides"],
        "warnings": result["warnings"],
    }


@app.post("/convert")
async def convert_endpoint(
    session_id: str = Form(...),
    output_file: str = Form(...),
):
    """Convert the merged PPTX to PDF (on-demand only)."""
    inp = OUTPUT_DIR / output_file
    if not inp.exists():
        raise HTTPException(404, "Merged file not found. Merge first.")

    pdf_name = f"converted_{session_id}.pdf"
    pdf_path = OUTPUT_DIR / pdf_name

    result = convert_to_pdf(inp, pdf_path)
    if not result["success"]:
        raise HTTPException(422, result["error"])

    return {"status": "success", "pdf_file": pdf_name, "method": result["method"]}


@app.get("/download")
async def download_file(
    session_id: str = Query(...),
    file_type: str = Query(...),
):
    """Serve merged PPTX or converted PDF for download."""
    if file_type not in ("pptx", "pdf"):
        raise HTTPException(400, "file_type must be 'pptx' or 'pdf'.")

    name = (
        f"merged_{session_id}.pptx"
        if file_type == "pptx"
        else f"converted_{session_id}.pdf"
    )
    path = OUTPUT_DIR / name
    if not path.exists():
        raise HTTPException(404, "File not found or has expired.")

    mime = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        if file_type == "pptx"
        else "application/pdf"
    )
    return FileResponse(
        str(path),
        filename=f"Merged_Presentation.{file_type}",
        media_type=mime,
    )


@app.get("/health")
async def health():
    return {"status": "healthy", "ts": datetime.now().isoformat()}


@app.on_event("startup")
async def on_startup():
    cleanup_stale()
    logger.info("PPTX Merger ready")


# ──────────────────────────────────────────────────────────────────
# Entry point
# ──────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", "8000"))
    uvicorn.run("app:app", host="0.0.0.0", port=port, reload=True)